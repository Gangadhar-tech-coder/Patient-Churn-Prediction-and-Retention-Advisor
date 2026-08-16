import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette
    PRIMARY_PURPLE = RGBColor(124, 58, 237)   # #7c3aed
    SECONDARY_CYAN = RGBColor(8, 145, 178)    # #0891b2
    DARK_SLATE = RGBColor(15, 23, 42)         # #0f172a
    TEXT_MUTED = RGBColor(100, 116, 139)      # #64748b
    BG_LIGHT = RGBColor(248, 250, 252)        # #f8fafc
    WHITE = RGBColor(255, 255, 255)
    RED_ACCENT = RGBColor(220, 38, 38)
    AMBER_ACCENT = RGBColor(217, 119, 6)
    GREEN_ACCENT = RGBColor(5, 150, 105)
    BORDER_COLOR = RGBColor(226, 232, 240)

    blank_layout = prs.slide_layouts[6]

    def add_header(slide, title_text, category_text="PATIENT CHURN & RETENTION ADVISOR"):
        # Header background banner
        header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.1))
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = DARK_SLATE
        header_shape.line.fill.background()

        # Accent top bar
        accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.08))
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = PRIMARY_PURPLE
        accent_bar.line.fill.background()

        # Category text
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.18), Inches(11), Inches(0.3))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = category_text.upper()
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = SECONDARY_CYAN

        # Main Title text
        txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(0.45), Inches(11), Inches(0.55))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = title_text
        p2.font.size = Pt(22)
        p2.font.bold = True
        p2.font.color.rgb = WHITE

    def add_card(slide, left, top, width, height, title="", bg_rgb=WHITE, border_rgb=BORDER_COLOR):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_rgb
        card.line.color.rgb = border_rgb
        card.line.width = Pt(1)
        
        if title:
            tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.4))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = DARK_SLATE
        return card

    # ==========================================
    # SLIDE 1: Title Slide
    # ==========================================
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = DARK_SLATE
    bg1.line.fill.background()

    # Title box
    tb = slide1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(3.0))
    tf = tb.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "PATIENT CHURN PREDICTION & RETENTION ADVISOR"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    p2 = tf.add_paragraph()
    p2.text = "An End-to-End Dual Machine Learning Architecture & Executive Clinical Advisory System"
    p2.font.size = Pt(18)
    p2.font.color.rgb = SECONDARY_CYAN
    p2.space_before = Pt(15)

    p3 = tf.add_paragraph()
    p3.text = "Predicting Individual & Cohort Churn Risk %, Diagnosing Root Causes & Prescribing 1-to-1 Retention Strategies"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(203, 213, 225)
    p3.space_before = Pt(10)

    # Metadata Card
    add_card(slide1, Inches(1.0), Inches(5.2), Inches(11.333), Inches(1.5), "", bg_rgb=RGBColor(30, 41, 59), border_rgb=PRIMARY_PURPLE)
    tb_meta = slide1.shapes.add_textbox(Inches(1.2), Inches(5.35), Inches(10.9), Inches(1.2))
    tf_meta = tb_meta.text_frame
    p_m = tf_meta.paragraphs[0]
    p_m.text = "PROJECT HIGHLIGHTS & ARCHITECTURE METADATA"
    p_m.font.size = Pt(11)
    p_m.font.bold = True
    p_m.font.color.rgb = SECONDARY_CYAN

    p_m2 = tf_meta.add_paragraph()
    p_m2.text = "• Dataset: 2,000 Enriched Patient Records  |  • ML Algorithms: Dual Random Forests (Binary + Multi-Class)  |  • ROC-AUC: 0.6065\n• Backend: FastAPI (Python 3.10) with Vectorized Batch Scoring (<50ms)  |  • Frontend: React 19, Vite 8, Tailwind CSS v4"
    p_m2.font.size = Pt(13)
    p_m2.font.color.rgb = WHITE
    p_m2.space_before = Pt(6)

    # ==========================================
    # SLIDE 2: Problem Statement & Objectives
    # ==========================================
    slide2 = prs.slides.add_slide(blank_layout)
    add_header(slide2, "1. Problem Statement & Executive Objectives")

    # Left Card - The Challenge
    add_card(slide2, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.4), "The Healthcare Churn Challenge")
    tb = slide2.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.2), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    bullets1 = [
        ("High Patient Attrition Risk", "Healthcare organizations lose up to 25% of their patient base annually due to unmanaged dissatisfaction and scheduling friction."),
        ("Silent Care Disengagement", "Patients quietly lapse without giving feedback, leading to missed preventive care, worsening chronic conditions, and lost revenue."),
        ("One-Size-Fits-All Retention", "Generic reminders fail because they do not address the specific root cause (e.g. wait times vs. billing issues)."),
        ("Computational Bottlenecks", "Legacy prediction systems process cohort data too slowly, delaying timely clinical outreach.")
    ]
    for title, desc in bullets1:
        p = tf.add_paragraph()
        p.text = f"• {title}: "
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_SLATE
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.size = Pt(12)
        run.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(12)

    # Right Card - Solution Objectives
    add_card(slide2, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.4), "Application Solution Objectives")
    tb2 = slide2.shapes.add_textbox(Inches(7.0), Inches(2.1), Inches(5.3), Inches(4.5))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    bullets2 = [
        ("Predict Exact Churn %", "Calculate individual and cohort churn probability (0–100%) and stratify into High, Medium, and Low risk tiers."),
        ("Diagnose Root Causes", "Pinpoint the exact primary driver of disengagement (e.g. missed appointments, billing friction, facility distance)."),
        ("Prescribe 1-to-1 Action Plans", "Map root causes directly to tailored clinical & administrative retention advice."),
        ("Vectorized Cohort Scoring", "Score 2,000+ patient records in <50ms to enable real-time population health management."),
        ("Interactive Re-engagement CTAs", "Empower care teams to trigger direct outreach calls, SMS rebooking, or advocate assignments in one click.")
    ]
    for title, desc in bullets2:
        p = tf2.add_paragraph()
        p.text = f"• {title}: "
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = PRIMARY_PURPLE
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.size = Pt(12)
        run.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(10)

    # ==========================================
    # SLIDE 3: System Architecture
    # ==========================================
    slide3 = prs.slides.add_slide(blank_layout)
    add_header(slide3, "2. End-to-End System Architecture")

    # 4 Layer Column Layout
    layers = [
        ("1. CLIENT PRESENTATION LAYER", "React 19 + Vite 8 + Tailwind CSS", [
            "MedVault Clinical Light Theme UI",
            "Collapsible 5-Section Sidebar",
            "Unified Risk Summary Hero Card",
            "Interactive Retention Strategy CTAs",
            "Real-Time Action Toast Alerts"
        ], PRIMARY_PURPLE),
        ("2. REST API & ROUTING LAYER", "FastAPI + Pydantic (Python 3.10)", [
            "CORS Middleware Enabled",
            "Single Patient Endpoint (/api/predict)",
            "Fast Cohort Endpoint (/api/batch-predict)",
            "Local & Render Cloud Dual Fallback",
            "Strict Pydantic Input Schemas"
        ], SECONDARY_CYAN),
        ("3. ML INFERENCE ENGINE", "Scikit-Learn + Joblib Artifacts", [
            "Vectorized Matrix Preprocessing",
            "Binary Churn Classifier (.predict_proba)",
            "Multi-Class Reason Classifier",
            "Clinical Feature Engineering",
            "1-to-1 Retention Strategy Mapper"
        ], DARK_SLATE),
        ("4. DATA & METRIC ENGINE", "Pandas + NumPy + Enriched Dataset", [
            "2,000 Patient Enriched CSV Data",
            "Feature Engineering Transformations",
            "Calculated Engagement Scores",
            "Calculated Cost-Per-Visit Ratios",
            "Satisfaction Index Aggregation"
        ], GREEN_ACCENT)
    ]

    col_width = Inches(2.7)
    gap = Inches(0.25)
    left_margin = Inches(0.8)

    for idx, (title, sub, items, color) in enumerate(layers):
        c_left = left_margin + idx * (col_width + gap)
        add_card(slide3, c_left, Inches(1.5), col_width, Inches(5.4), "", bg_rgb=WHITE)
        
        # Header banner inside column
        banner = slide3.shapes.add_shape(MSO_SHAPE.RECTANGLE, c_left, Inches(1.5), col_width, Inches(0.8))
        banner.fill.solid()
        banner.fill.fore_color.rgb = color
        banner.line.fill.background()

        tb = slide3.shapes.add_textbox(c_left + Inches(0.1), Inches(1.55), col_width - Inches(0.2), Inches(0.7))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = WHITE

        p_s = tf.add_paragraph()
        p_s.text = sub
        p_s.font.size = Pt(9)
        p_s.font.color.rgb = RGBColor(241, 245, 249)

        # Content list
        tb_c = slide3.shapes.add_textbox(c_left + Inches(0.15), Inches(2.4), col_width - Inches(0.3), Inches(4.3))
        tf_c = tb_c.text_frame
        tf_c.word_wrap = True
        for item in items:
            p_i = tf_c.add_paragraph()
            p_i.text = f"✔ {item}"
            p_i.font.size = Pt(11)
            p_i.font.color.rgb = DARK_SLATE
            p_i.space_after = Pt(10)

    # ==========================================
    # SLIDE 4: Technology Stack
    # ==========================================
    slide4 = prs.slides.add_slide(blank_layout)
    add_header(slide4, "3. Technology Stack & Frameworks")

    stacks = [
        ("Frontend Technology Stack", [
            ("Framework & Library", "React 19 & Vite 8 (Ultra-fast HMR & production build)"),
            ("Styling System", "Tailwind CSS v4 + MedVault Vanilla CSS Design Tokens"),
            ("Typography", "Google Fonts Inter (Modern medical analytics typography)"),
            ("Visual Components", "Custom SVG Gauge Chart, Metric Cards, Collapsible Sidebars"),
            ("State & API Layer", "React Hooks (useState, useEffect) + Native Fetch API with Dual Fallback")
        ]),
        ("Backend & Infrastructure Stack", [
            ("API Framework", "FastAPI (High performance, async native, auto OpenAPI)"),
            ("Server Engine", "Uvicorn ASGI Server (Port 8000 local dev & production)"),
            ("Data Validation", "Pydantic Schemas (Strict input validation & response typing)"),
            ("Machine Learning", "Scikit-Learn, Pandas, NumPy, Joblib"),
            ("Version Control & Cloud", "Git, GitHub Remote Repository, Render Cloud Platform")
        ])
    ]

    for idx, (title, items) in enumerate(stacks):
        left = Inches(0.8 + idx * 5.9)
        add_card(slide4, left, Inches(1.5), Inches(5.6), Inches(5.4), title)
        tb = slide4.shapes.add_textbox(left + Inches(0.2), Inches(2.1), Inches(5.2), Inches(4.5))
        tf = tb.text_frame
        tf.word_wrap = True
        for label, desc in items:
            p = tf.add_paragraph()
            p.text = f"• {label}: "
            p.font.bold = True
            p.font.size = Pt(13)
            p.font.color.rgb = PRIMARY_PURPLE if idx == 0 else SECONDARY_CYAN
            run = p.add_run()
            run.text = desc
            run.font.bold = False
            run.font.size = Pt(12)
            run.font.color.rgb = TEXT_MUTED
            p.space_after = Pt(14)

    # ==========================================
    # SLIDE 5: Feature Engineering & Data Signals
    # ==========================================
    slide5 = prs.slides.add_slide(blank_layout)
    add_header(slide5, "4. Data Signals & Feature Engineering Engine")

    # Left Box - Input Features
    add_card(slide5, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.4), "Input Feature Spectrum (21 Signals)")
    tb = slide5.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.2), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    features = [
        ("Demographics", "Age, Gender, State"),
        ("Clinical Attributes", "Specialty, Insurance Type, Tenure (Months)"),
        ("Engagement Metrics", "Visits Last Year, Missed Appointments, Days Since Last Visit, Portal Usage, Referrals Made"),
        ("Satisfaction Scores", "Overall Satisfaction, Wait Time Satisfaction, Staff Satisfaction, Provider Rating"),
        ("Financial & Access", "Avg Out-of-Pocket Cost, Billing Issues, Distance to Facility (Miles)")
    ]
    for cat, items in features:
        p = tf.add_paragraph()
        p.text = f"• {cat}: "
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_SLATE
        run = p.add_run()
        run.text = items
        run.font.bold = False
        run.font.size = Pt(12)
        run.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(10)

    # Right Box - Engineered Formulae
    add_card(slide5, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.4), "Domain Feature Engineering Metrics")
    tb2 = slide5.shapes.add_textbox(Inches(7.0), Inches(2.1), Inches(5.3), Inches(4.5))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    engineered = [
        ("1. Engagement Score", "Formula: Visits_Last_Year - Missed_Appointments", "Measures true active healthcare utilization versus missed commitment rate."),
        ("2. Cost Per Visit", "Formula: Avg_Out_Of_Pocket_Cost / (Visits_Last_Year + 1)", "Quantifies financial burden per encounter to detect cost sensitivity."),
        ("3. Satisfaction Index Avg", "Formula: (Overall + Wait_Time + Staff) / 3", "Blends holistic patient satisfaction into a single composite score."),
        ("4. One-Hot Categorical Encoding", "Encoded Features: Gender, State, Specialty, Insurance", "Transforms 4 categorical variables into a clean 48-column model matrix.")
    ]
    for title, formula, desc in engineered:
        p = tf2.add_paragraph()
        p.text = f"• {title}"
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = SECONDARY_CYAN
        
        p_f = tf2.add_paragraph()
        p_f.text = f"  {formula}"
        p_f.font.bold = True
        p_f.font.size = Pt(11)
        p_f.font.color.rgb = PRIMARY_PURPLE

        p_d = tf2.add_paragraph()
        p_d.text = f"  {desc}"
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = TEXT_MUTED
        p_d.space_after = Pt(8)

    # ==========================================
    # SLIDE 6: Model 1 Breakdown
    # ==========================================
    slide6 = prs.slides.add_slide(blank_layout)
    add_header(slide6, "5. Model 1: Binary Churn Probability Classifier")

    # Left Column - Hyperparameters & Training
    add_card(slide6, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.4), "Model Specifications & Training")
    tb = slide6.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.2), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    specs = [
        ("Algorithm", "RandomForestClassifier (Scikit-Learn Ensemble)"),
        ("Trees (n_estimators)", "300 Decision Trees"),
        ("Max Tree Depth", "12 (Prevents overfitting while capturing non-linear interactions)"),
        ("Min Samples Split", "4 samples"),
        ("Train/Test Split", "80% Training (1,600 rows), 20% Testing (400 rows) with Stratified Sampling"),
        ("Output Probability", "predict_proba[:, 1] yields continuous probability from 0.0000 to 1.0000 (0% to 100%)")
    ]
    for label, val in specs:
        p = tf.add_paragraph()
        p.text = f"• {label}: "
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_SLATE
        run = p.add_run()
        run.text = val
        run.font.bold = False
        run.font.size = Pt(12)
        run.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(10)

    # Right Column - Risk Stratification Tiers
    add_card(slide6, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.4), "Risk Stratification Thresholds")
    
    tiers = [
        ("HIGH CHURN RISK 🔴", "Probability >= 65.0%", "Immediate clinical intervention required. High likelihood of patient disengagement within 60 days.", RED_ACCENT, RGBColor(254, 242, 242)),
        ("MEDIUM CHURN RISK 🟡", "Probability 45.0% - 64.9%", "Proactive retention monitoring recommended. Patient shows early warning disengagement signals.", AMBER_ACCENT, RGBColor(255, 251, 235)),
        ("LOW CHURN RISK 🟢", "Probability < 45.0%", "Patient is well-engaged and satisfied with current care services. Maintain routine protocol.", GREEN_ACCENT, RGBColor(236, 253, 245))
    ]

    for idx, (title, range_t, desc, color, bg) in enumerate(tiers):
        t_top = Inches(2.2 + idx * 1.5)
        add_card(slide6, Inches(7.0), t_top, Inches(5.3), Inches(1.35), "", bg_rgb=bg, border_rgb=color)
        tb_t = slide6.shapes.add_textbox(Inches(7.2), t_top + Inches(0.1), Inches(4.9), Inches(1.15))
        tf_t = tb_t.text_frame
        tf_t.word_wrap = True
        
        p = tf_t.paragraphs[0]
        p.text = f"{title} ({range_t})"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = color
        
        p_d = tf_t.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(11)
        p_d.font.color.rgb = DARK_SLATE
        p_d.space_before = Pt(4)

    # ==========================================
    # SLIDE 7: Model 2 Breakdown
    # ==========================================
    slide7 = prs.slides.add_slide(blank_layout)
    add_header(slide7, "6. Model 2: Multi-Class Churn Reason & Advice Mapper")

    add_card(slide7, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.4), "Multi-Class Reason Model (`reason_model.pkl`)")
    tb = slide7.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.2), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    m2_specs = [
        ("Algorithm", "Multi-Class RandomForestClassifier + LabelEncoder"),
        ("N Estimators & Depth", "300 Trees, max_depth=14"),
        ("Primary Function", "Diagnoses the root cause reason why a patient disengages when risk >= 45%."),
        ("Diagnosed Categories", "Long gap since last visit, missed appointment patterns, billing friction, out-of-pocket cost burden, wait times, staff experience, facility distance.")
    ]
    for label, val in m2_specs:
        p = tf.add_paragraph()
        p.text = f"• {label}: "
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = SECONDARY_CYAN
        run = p.add_run()
        run.text = val
        run.font.bold = False
        run.font.size = Pt(12)
        run.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(12)

    add_card(slide7, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.4), "1-to-1 Retention Strategy Mapping (`advice_map.pkl`)")
    tb2 = slide7.shapes.add_textbox(Inches(7.0), Inches(2.1), Inches(5.3), Inches(4.5))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    mappings = [
        ("Lapsed Patient (>250 Days)", "Send re-engagement outreach with an easy rebooking link & checkup reminder."),
        ("Frequently Missed Appts", "Discuss schedule barriers & offer virtual / telehealth options."),
        ("Billing Issues / High Cost", "Connect patient with financial counseling & flexible payment plans."),
        ("Low Satisfaction / Provider", "Assign patient advocate & schedule quality care follow-up call.")
    ]
    for cause, advice in mappings:
        p = tf2.add_paragraph()
        p.text = f"• Diagnosed Cause: {cause}"
        p.font.bold = True
        p.font.size = Pt(12)
        p.font.color.rgb = DARK_SLATE

        p_a = tf2.add_paragraph()
        p_a.text = f"  ➡ Retention Strategy: {advice}"
        p_a.font.size = Pt(11)
        p_a.font.color.rgb = PRIMARY_PURPLE
        p_a.space_after = Pt(10)

    # ==========================================
    # SLIDE 8: Evaluation & Performance
    # ==========================================
    slide8 = prs.slides.add_slide(blank_layout)
    add_header(slide8, "7. Model Metrics & Performance Benchmarks")

    # Metric Cards Top Row
    m_boxes = [
        ("BINARY MODEL ROC-AUC", "0.6065", "Evaluated on 400 Stratified Holdout Test Patients", PRIMARY_PURPLE),
        ("REASON MODEL ACCURACY", "High", "Multi-Class Label Classification Accuracy", SECONDARY_CYAN),
        ("BATCH INFERENCE SPEED", "<50ms", "Vectorized Prediction for 2,000 Patient Cohort", GREEN_ACCENT)
    ]
    for idx, (title, val, desc, color) in enumerate(m_boxes):
        left = Inches(0.8 + idx * 3.9)
        add_card(slide8, left, Inches(1.5), Inches(3.7), Inches(1.8), "", bg_rgb=WHITE, border_rgb=color)
        tb_m = slide8.shapes.add_textbox(left + Inches(0.15), Inches(1.65), Inches(3.4), Inches(1.5))
        tf_m = tb_m.text_frame
        tf_m.word_wrap = True
        p = tf_m.paragraphs[0]
        p.text = title
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = TEXT_MUTED

        p_v = tf_m.add_paragraph()
        p_v.text = val
        p_v.font.size = Pt(28)
        p_v.font.bold = True
        p_v.font.color.rgb = color

        p_d = tf_m.add_paragraph()
        p_d.text = desc
        p_d.font.size = Pt(9)
        p_d.font.color.rgb = DARK_SLATE

    # Optimization Comparison Bottom Box
    add_card(slide8, Inches(0.8), Inches(3.6), Inches(11.7), Inches(3.3), "Batch Prediction Optimization Benchmark (2,000 Records)")
    tb = slide8.shapes.add_textbox(Inches(1.0), Inches(4.2), Inches(11.3), Inches(2.5))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.add_paragraph()
    p.text = "• Legacy Row-by-Row Iteration (Unoptimized Loop):"
    p.font.bold = True
    p.font.size = Pt(13)
    p.font.color.rgb = RED_ACCENT
    run = p.add_run()
    run.text = " Previously iterated over df.iterrows(), constructing Pydantic models 2,000 times. Resulted in 3.2s – 5.5s execution latency."
    run.font.bold = False
    run.font.size = Pt(12)
    run.font.color.rgb = TEXT_MUTED
    p.space_after = Pt(14)

    p2 = tf.add_paragraph()
    p2.text = "• Vectorized Batch Prediction Engine (Optimized predict_batch):"
    p2.font.bold = True
    p2.font.size = Pt(13)
    p2.font.color.rgb = GREEN_ACCENT
    run2 = p2.add_run()
    run2.text = " Replaced row loops with C-optimized Pandas/Scikit-Learn matrix transformations. Processes 2,000 rows in <50ms (>100x Speedup)."
    run2.font.bold = False
    run2.font.size = Pt(12)
    run2.font.color.rgb = TEXT_MUTED

    # ==========================================
    # SLIDE 9: Application Functionalities
    # ==========================================
    slide9 = prs.slides.add_slide(blank_layout)
    add_header(slide9, "8. Application Features & UI Capabilities")

    funcs = [
        ("Single Patient Profiler", "Interactive sidebar with 5 collapsible sections & range sliders to score individual patients in real-time.", PRIMARY_PURPLE),
        ("Cohort Batch Predictor", "Drag-and-drop CSV upload for cohort-wide screening with automatic risk summaries & data tables.", SECONDARY_CYAN),
        ("Unified Risk Hero Card", "Combines SVG percentage gauge, risk tier badge, and Diagnosed Root Cause panel into one hero view.", DARK_SLATE),
        ("Direct Action CTAs", "Direct call-to-action buttons (Outreach Call, SMS Rebooking, Advocate Assignment, Telehealth) with Toast alerts.", GREEN_ACCENT)
    ]

    for idx, (title, desc, color) in enumerate(funcs):
        left = Inches(0.8 + (idx % 2) * 5.9)
        top = Inches(1.5 + (idx // 2) * 2.7)
        add_card(slide9, left, top, Inches(5.6), Inches(2.4), title, border_rgb=color)
        tb = slide9.shapes.add_textbox(left + Inches(0.2), top + Inches(0.65), Inches(5.2), Inches(1.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_SLATE

    # ==========================================
    # SLIDE 10: Business Impact & Roadmap
    # ==========================================
    slide10 = prs.slides.add_slide(blank_layout)
    add_header(slide10, "9. Business Impact & Future Roadmap")

    add_card(slide10, Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.4), "Clinical & Business Impact")
    tb = slide10.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.2), Inches(4.5))
    tf = tb.text_frame
    tf.word_wrap = True
    impacts = [
        ("Reduces Patient Attrition", "Early detection allows care teams to intervene before patients permanently disengage."),
        ("Improves Care Continuity", "Ensures high-risk patients receive timely preventive care and chronic disease management."),
        ("Protects Health System Revenue", "Retaining existing patients preserves lifetime patient value and reduces acquisition costs."),
        ("Actionable Workflows", "Replaces static reports with actionable outreach triggers directly integrated into clinical workflows.")
    ]
    for title, desc in impacts:
        p = tf.add_paragraph()
        p.text = f"• {title}: "
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = DARK_SLATE
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.size = Pt(12)
        run.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(12)

    add_card(slide10, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.4), "Future Technical Roadmap")
    tb2 = slide10.shapes.add_textbox(Inches(7.0), Inches(2.1), Inches(5.3), Inches(4.5))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    roadmap = [
        ("EHR / FHIR Integration", "Connect directly to Epic / Cerner via FHIR APIs for automated real-time patient syncing."),
        ("Automated SMS Gateway", "Integrate Twilio / AWS SNS to trigger automated SMS rebooking messages upon CTA click."),
        ("Longitudinal Survival Analysis", "Incorporate Kaplan-Meier curves to model time-to-churn dynamics."),
        ("Advanced SHAP Explainability", "Add SHAP (SHapley Additive exPlanations) values for deep feature attribution per patient.")
    ]
    for title, desc in roadmap:
        p = tf2.add_paragraph()
        p.text = f"• {title}: "
        p.font.bold = True
        p.font.size = Pt(13)
        p.font.color.rgb = PRIMARY_PURPLE
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.size = Pt(12)
        run.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(12)

    output_path = os.path.join("d:\\Desktop\\cts2\\Patient churn predictor 2", "Patient_Churn_Prediction_and_Retention_Advisor_Presentation.pptx")
    prs.save(output_path)
    print(f"Presentation saved successfully to {output_path}")

if __name__ == "__main__":
    create_presentation()
